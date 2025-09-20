"""
Document Processor for Local Mind
Handles document parsing, chunking, and position tracking
"""

import os
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
import json
import logging
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
import markdown

from config.app_settings import CHUNK_SIZE, CHUNK_OVERLAP, DATA_DIR

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Process documents and extract content with position tracking"""

    def __init__(self, chunk_size: int = CHUNK_SIZE, chunk_overlap: int = CHUNK_OVERLAP):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx
        }

    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        Process a file and extract content with metadata

        Args:
            file_path: Path to the file to process

        Returns:
            Dictionary with file metadata and chunks
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        # Get file extension
        ext = path.suffix.lower()
        if ext not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {ext}")

        # Generate document ID
        doc_id = self._generate_document_id(file_path)

        # Process file based on type
        processor = self.supported_formats[ext]
        chunks = processor(file_path)

        # Add metadata to each chunk
        for i, chunk in enumerate(chunks):
            chunk['document_id'] = doc_id
            chunk['chunk_id'] = f"{doc_id}_chunk_{i}"
            chunk['file_path'] = str(path.absolute())
            chunk['file_name'] = path.name
            chunk['chunk_index'] = i
            chunk['total_chunks'] = len(chunks)

        return {
            'document_id': doc_id,
            'file_path': str(path.absolute()),
            'file_name': path.name,
            'file_size': path.stat().st_size,
            'file_type': ext,
            'chunks': chunks,
            'total_chunks': len(chunks),
            'processed_at': datetime.now().isoformat()
        }

    def _generate_document_id(self, file_path: str) -> str:
        """Generate a unique document ID based on file path and content"""
        path = Path(file_path)
        content_hash = hashlib.md5(
            f"{path.absolute()}_{path.stat().st_mtime}".encode()
        ).hexdigest()[:12]
        return f"doc_{content_hash}"

    def _process_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PDF with position tracking using PyMuPDF"""
        chunks = []

        try:
            doc = fitz.open(file_path)

            for page_num, page in enumerate(doc):
                # Extract text blocks with position
                blocks = page.get_text("blocks")

                page_text = ""
                page_positions = []

                for block in blocks:
                    # block format: (x0, y0, x1, y1, "text", block_no, block_type)
                    if len(block) >= 5:
                        bbox = {
                            "x0": block[0],
                            "y0": block[1],
                            "x1": block[2],
                            "y1": block[3]
                        }
                        text = block[4].strip()

                        if text:
                            page_text += text + "\n"
                            page_positions.append({
                                "text": text,
                                "bbox": bbox,
                                "page": page_num + 1
                            })

                # Chunk the page text while preserving positions
                if page_text.strip():
                    page_chunks = self._chunk_text_with_positions(
                        page_text,
                        page_positions,
                        page_num + 1
                    )
                    chunks.extend(page_chunks)

            doc.close()

        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {e}")
            raise

        return chunks

    def _process_text(self, file_path: str) -> List[Dict[str, Any]]:
        """Process plain text file"""
        chunks = []

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Simple chunking for text files
            text_chunks = self._simple_chunk_text(content)

            for i, chunk_text in enumerate(text_chunks):
                chunks.append({
                    'content': chunk_text,
                    'page': 0,
                    'bbox': None,
                    'metadata': {
                        'source': 'text_file',
                        'chunk_method': 'simple'
                    }
                })

        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
            raise

        return chunks

    def _process_markdown(self, file_path: str) -> List[Dict[str, Any]]:
        """Process markdown file"""
        chunks = []

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            # Convert markdown to HTML for structure
            html_content = markdown.markdown(content)

            # Extract headers and sections
            sections = self._extract_markdown_sections(content)

            for section in sections:
                # Chunk each section
                text_chunks = self._simple_chunk_text(section['content'])

                for chunk_text in text_chunks:
                    chunks.append({
                        'content': chunk_text,
                        'page': 0,
                        'bbox': None,
                        'metadata': {
                            'source': 'markdown_file',
                            'section': section.get('title', ''),
                            'level': section.get('level', 0)
                        }
                    })

        except Exception as e:
            logger.error(f"Error processing markdown file {file_path}: {e}")
            raise

        return chunks

    def _process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        """Process Word document"""
        chunks = []

        try:
            doc = DocxDocument(file_path)
            full_text = []

            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)

            # Join and chunk the text
            content = "\n".join(full_text)
            text_chunks = self._simple_chunk_text(content)

            for i, chunk_text in enumerate(text_chunks):
                chunks.append({
                    'content': chunk_text,
                    'page': 0,  # Word docs don't have pages in the same way
                    'bbox': None,
                    'metadata': {
                        'source': 'docx_file',
                        'paragraph_count': len(full_text)
                    }
                })

        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")
            raise

        return chunks

    def _process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PowerPoint presentation"""
        chunks = []

        try:
            prs = Presentation(file_path)

            for slide_num, slide in enumerate(prs.slides):
                slide_text = []

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    slide_content = "\n".join(slide_text)
                    chunks.append({
                        'content': slide_content,
                        'page': slide_num + 1,
                        'bbox': None,
                        'metadata': {
                            'source': 'pptx_file',
                            'slide_number': slide_num + 1,
                            'slide_title': slide_text[0] if slide_text else ""
                        }
                    })

        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")
            raise

        return chunks

    def _chunk_text_with_positions(
        self,
        text: str,
        positions: List[Dict[str, Any]],
        page_num: int
    ) -> List[Dict[str, Any]]:
        """Chunk text while preserving position information"""
        chunks = []
        words = text.split()

        if not words:
            return chunks

        # Simple word-based chunking
        for i in range(0, len(words), self.chunk_size - self.chunk_overlap):
            chunk_words = words[i:i + self.chunk_size]
            chunk_text = " ".join(chunk_words)

            # Find the first and last position for this chunk
            # This is simplified - in production, you'd want more precise mapping
            bbox = None
            if positions:
                bbox = positions[0]['bbox']  # Simplified: use first block's bbox

            chunks.append({
                'content': chunk_text,
                'page': page_num,
                'bbox': bbox,
                'metadata': {
                    'chunk_method': 'position_aware',
                    'word_count': len(chunk_words)
                }
            })

        return chunks

    def _simple_chunk_text(self, text: str) -> List[str]:
        """Simple text chunking by character count with overlap"""
        chunks = []
        text_length = len(text)

        if text_length <= self.chunk_size:
            return [text]

        for i in range(0, text_length, self.chunk_size - self.chunk_overlap):
            chunk = text[i:i + self.chunk_size]
            chunks.append(chunk)

            # If we've covered all text, break
            if i + self.chunk_size >= text_length:
                break

        return chunks

    def _extract_markdown_sections(self, content: str) -> List[Dict[str, Any]]:
        """Extract sections from markdown content"""
        sections = []
        lines = content.split('\n')
        current_section = {'title': 'Introduction', 'level': 0, 'content': []}

        for line in lines:
            if line.startswith('#'):
                # Save current section if it has content
                if current_section['content']:
                    current_section['content'] = '\n'.join(current_section['content'])
                    sections.append(current_section)

                # Start new section
                level = len(line) - len(line.lstrip('#'))
                title = line.lstrip('#').strip()
                current_section = {'title': title, 'level': level, 'content': []}
            else:
                current_section['content'].append(line)

        # Add the last section
        if current_section['content']:
            current_section['content'] = '\n'.join(current_section['content'])
            sections.append(current_section)

        return sections if sections else [{'title': 'Document', 'level': 0, 'content': content}]

    def process_folder(self, folder_path: str, recursive: bool = True) -> List[Dict[str, Any]]:
        """
        Process all supported files in a folder

        Args:
            folder_path: Path to the folder
            recursive: Whether to process subfolders

        Returns:
            List of processed document metadata
        """
        processed_docs = []
        folder = Path(folder_path)

        if not folder.exists() or not folder.is_dir():
            raise ValueError(f"Invalid folder path: {folder_path}")

        # Get all supported files
        pattern = '**/*' if recursive else '*'
        for file_path in folder.glob(pattern):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_formats:
                try:
                    doc_data = self.process_file(str(file_path))
                    processed_docs.append(doc_data)
                    logger.info(f"Processed: {file_path.name}")
                except Exception as e:
                    logger.error(f"Failed to process {file_path}: {e}")
                    continue

        return processed_docs